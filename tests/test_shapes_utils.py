########################################################################################################################
# Filename: test_shape_utils.py
#
# Purpose: Tests for shapes_utils.
# Author(s): Bobby (Robert) Lumpkin
#
# Library Dependencies: 
########################################################################################################################

import pytest
import pptx
from copy import deepcopy
import pandas as pd
import io
import sys

sys.path.append("C:\\Users\\rober\\OneDrive\\Documents\\Github_Repos\\pptx_utils\\src")
from shape_utils import init_pres, table_from_df, matplotlib_to_pic


@pytest.mark.shape_utils
def test_init_pres():
    """
    Tests the functionality of 'init_pres()'.
    """
    ppt_stream = io.BytesIO()
    # Call init_pres()
    ppt = init_pres(
        title = "Test Presentation",
        template = None,
        save = True,
        save_path = ppt_stream
    )
    
    # Check save functionality
    assert ppt_stream.getbuffer().nbytes > 0
    loaded_ppt = pptx.Presentation(ppt_stream)

    # Check slide title
    try:
        assert ppt.slides[0].shapes.title.text == "Test Presentation"
        assert loaded_ppt.slides[0].shapes.title.text == "Test Presentation"
    except IndexError:
        assert False


@pytest.mark.shape_utils
def test_table_from_df(
    slide,
    df,
    col_rename_dict
):
    """
    Tets functionality of 'table_from_df()'.
    """
    args_list = [{"round_floats" : True, 
                  "title" : None},
                 {"round_floats" : False,
                  "title" : "Test Table"}]
    slides = {"title" : slide,
              "no_title" : deepcopy(slide)}
    for args in args_list:
        # Call table_from_df()
        if args["title"] is None:
            args_key = "no_title"
        else:
            args_key = "title"
        new_slide = table_from_df(
            slide = slides[args_key],
            df = df,
            top = 2,
            left = 1,
            width = 2,
            height = 2,
            font_size = 12,
            round_digits = 2,
            header_font_size = 14,
            title_font_size = 18, 
            col_rename_dict = col_rename_dict, 
            **args
        )

        # Check that a table exists, is of the correct size and that columns were renamed
        assert new_slide.shapes[1].has_table
        num_cols = len(new_slide.shapes[1].table.columns)
        num_rows = len(new_slide.shapes[1].table.rows)
        df_cols = df.columns.tolist()
        if args["title"] is None:
            assert new_slide.shapes[1].table.cell(0, 0).text == "NEW_" + df_cols[0]
            assert new_slide.shapes[1].table.cell(0, 1).text == "NEW_" + df_cols[1]
            assert num_cols == df.shape[1]
            assert num_rows == df.shape[0] + 1
        else:
            assert new_slide.shapes[1].table.cell(0, 0).text == "Test Table"
            assert new_slide.shapes[1].table.cell(1, 0).text == "NEW_" + df_cols[0]
            assert new_slide.shapes[1].table.cell(1, 1).text == "NEW_" + df_cols[1]
            assert num_cols == df.shape[1] 
            assert num_rows == df.shape[0] + 2
        
    
    