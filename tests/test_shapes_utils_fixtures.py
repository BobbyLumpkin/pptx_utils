########################################################################################################################
# Filename: test_shape_utils_fixtures.py
#
# Purpose: Fixtures for test_shapes_utils.
# Author(s): Bobby (Robert) Lumpkin
#
# Library Dependencies: 
########################################################################################################################

import pytest
import pptx
import pandas as pd
from random import randrange


@pytest.fixture()
def slide():
    ppt = pptx.Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = "Test Slide"
    return slide
    

@pytest.fixture()
def df():
    return pd.DataFrame({"Col1" : [1,2,3], 
                         "Col2" : ["Hello", "World", "!"]})


@pytest.fixture()
def col_rename_dict(
    df
):
    col_names = df.columns.tolist()
    rename_tuples = [(col_name, "NEW_" + col_name) for col_name in col_names]
    return dict(rename_tuples)

