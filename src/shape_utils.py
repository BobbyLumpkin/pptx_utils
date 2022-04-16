########################################################################################################################
# Filename: shape_utils.py
#
# Purpose: Utilities for pptx.shapes.
# Author(s): Bobby (Robert) Lumpkin
#
# Library Dependencies: 
########################################################################################################################

import pandas as pd
import numpy as np
from numbers import Number
import io
import pptx
from pptx.util import Inches, Pt

def init_pres(
    title: str,
    template: str = None,
    save: bool = False,
    save_path: str = None
):
    """
    Initializes a pptx.Presentation object.
    """
    ppt = pptx.Presentation()
    if template:
        ppt = pptx.Presentation(template)
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = title

    # Save, if appropriate
    if save:
        ppt.save(save_path)

    return ppt


def table_from_d(
    slide,
    df: pd.DataFrame,
    top: float,
    left: float,
    width: float,
    height: float,
    font_size: int = 12,
    round_floats: bool = True,
    round_digits: int = 2,
    header_font_size: int = 14,
    title: str = None,
    title_font_size: int = 18, 
    col_rename_dict: dict = None
):
    """
    Converts a pandas dataframe into a table and adds to a pptx.Presentation slide.

    Parameters
    ----------
    slide: the slide object to add the table to
    df: the pandas data frame to convert into a table
    top: integer distance of the top edge of this shape from the top edge of the slide
    left: integer distance of the left edge of this shape from the left edge of the slide
    width: integer distance between left and right extents of shape
    height: integer distance between top and bottom extents of shape
    font_size: the font size for table entries. Defaults to 12 pt.
    round_floats: whether or not to round float types. Defaults to True.
    round_digits: the number of decimals to use when rounding float types. Defaults to 2.
    header_font_size: the font size for table column headers. Defaults to 14 pt.
    title: the table title, to appear above the table (optional).
    title_font_size: the font size for the table title. Defaults to 18 pt.
    col_rename_dict: a dictionary with current df column names as keys and 
        desired table header names as values (optional).

    Returns
    ----------
    slide
    """
    num_rows = df.shape[0]
    num_cols = df.shape[1]

    # Initialize the table
    if title:
        num_rows = num_rows + 1
    shape = slide.shapes.add_table(
        rows = num_rows + 1,
        cols = num_cols,
        left = Inches(left),
        top = Inches(top),
        width = Inches(width),
        height = Inches(height)
    )
    table = shape.table

    # Add title and column names
    header_row = 0
    if title:
        cell = table.cell(0, 0)
        cell.merge(table.cell(0, num_cols - 1))
        cell.text = title
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(title_font_size)
        header_row = 1
    for i in list(range(0, num_cols)):
        col = df.columns.tolist()[i]
        cell = table.cell(header_row, i)
        cell.text = col
        if col_rename_dict:
            cell.text = col_rename_dict[col]
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(header_font_size)
        paragraph.font.bold = True
        paragraph.font.underline = True

    # Add df contents
    for i in list(range(header_row + 1, num_rows + 1)):
        for j in list(range(0, num_cols)):
            cell = table.cell(i, j)
            col_name = table.cell(header_row, j).text
            if col_rename_dict:
                rev_col_rename_dict = {v : k for k, v in col_rename_dict.items()}
                col_name = rev_col_rename_dict[col_name]
            iloc_val = i - 2
            if not title:
                iloc_val = i - 1
            cell_value = df[col_name].iloc[iloc_val]
            if isinstance(cell_value, Number):
                if round_floats:
                    cell_value = round(cell_value, round_digits)
                cell_value = "{:,}".format(cell_value)
            cell.text = str(cell_value)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(font_size)

    return slide


def matplotlib_to_pic(
    slide,
    figure,
    left: int,
    top: int,
    width: int,
    height: int
):
    """
    Adds a matplotlib figure to a slide as a picture.

    Parameters 
    ----------
    slide: a slide object to add the figure to 
    figure: a matplotlib figure to be included in 'slide'
    left: integer distance of the left edge of this shape from the left edge of the slide
    top: integer distance of the top edge of this shape from the top edge of this slied
    width: integer distance between left and right extents of shape
    height: integer distance between top and bottom extents of shape

    Returns
    ----------
    slide
    """
    image_stream = io.BytesIO()
    figure.savefig(image_stream)
    pic = slide.shapes.add_picture(
        image_file = image_stream,
        left = Inches(left),
        top = Inches(top),
        width = Inches(width),
        height = Inches(height)
    )

    return slide